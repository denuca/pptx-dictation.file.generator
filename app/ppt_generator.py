import logging
from core.file_generator.pptx import create_ppt, add_title_slide  # Import from the core library
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
from PIL import ImageFont
import re
import json
import os

# Set up logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

# Load styles from JSON file
try:
    with open('styles.json', 'r') as f:
        styles = json.load(f)
except Exception as e:
    logging.error(f"Failed to load styles from JSON: {e}")
    raise  # Re-raise the exception after logging

def create_ppt_from_text(title, content):
    """Create a PowerPoint presentation from a title and content."""
    logging.debug("Creating PowerPoint from text...")
    try:
        ppt = create_ppt()  # Create a new PowerPoint presentation
        add_title_slide(ppt, title, content)  # Add the title slide
        return ppt
    except Exception as e:
        logging.error(f"Error creating PowerPoint: {e}")
        raise  # Re-raise the exception after logging

def create_ppt_dictation_from_text(title, content):
    """Create a PowerPoint presentation for dictation from title and content."""
    logging.debug("Creating PowerPoint dictation from text...")
    try:
        ppt = create_ppt()  # Create a new PowerPoint presentation
        slide = add_content_slide(ppt, title, content)  # Add content slide
        logging.debug(f"Added content slide: {slide.slide_layout.name}")

        # Split the content into lines based on carriage returns
        lines = content.split('\n')
        
        max_height = 0
        for line in lines:
            words = re.findall(r'\w+|[^\w\s]', line, re.UNICODE)
            for word in words:
                content_type = 'word' if re.match(r'\w+', word) else 'punctuation'
                _, text_height = calculate_text_size(word, styles[content_type]['font_path'], styles[content_type]['font_size'])
                max_height = max(max_height, Inches(text_height / 72))

        # Initialize placeholder, slide number, and word ID
        content_placeholder = slide.placeholders[1]
        slide_number = 1
        word_id = 1
        left = content_placeholder.left
        top = content_placeholder.top
        
        for line in lines:
            words = re.findall(r'\w+|[^\w\s]', line, re.UNICODE)
            for word in words:
                content_type = 'word' if re.match(r'\w+', word) else 'punctuation'
                add_text_box(slide, word, left, top + Inches(0.3), max_height, word_id, content_type)  # Add text box
                text_width, _ = calculate_text_size(word, styles[content_type]['font_path'], styles[content_type]['font_size'])
                left += Inches(text_width / 72) + Inches(0.2)  # Adjust spacing between text boxes

                if styles[content_type]['count_id']:
                    word_id += 1

                # Check if the next text box will fit in the current line
                if left + Inches(1.5) > content_placeholder.left + content_placeholder.width:
                    # Move to the next line
                    left = content_placeholder.left  # Reset to placeholder's left
                    top += max_height + Inches(0.6)  # Move down to the next line
                    
                    # Check if the next text box will fit in the current slide
                    if top + max_height + Inches(0.6) > content_placeholder.top + content_placeholder.height:
                        # Add a new slide
                        slide_number += 1
                        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                        title_shape = slide.shapes.title
                        title_shape.text = title
                        content_placeholder = slide.placeholders[1]
                        top = content_placeholder.top

            # Move to the next line after processing all words in the current line
            left = content_placeholder.left  # Ensure the left is reset to align on each new line
            top += max_height + Inches(0.6)

        # Clean up content placeholders and add headers/footers as needed
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    sp = shape._element
                    sp.getparent().remove(sp)

        # Add headers and footers
        for i, slide in enumerate(ppt.slides):
            header = slide.shapes.add_textbox(Cm(1), Cm(0.5), ppt.slide_width - Cm(2), Cm(1))
            header_frame = header.text_frame
            header_frame.text = styles['header']['text']
            header_frame.paragraphs[0].font.size = Pt(styles['header']['font_size'])
            header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            footer = slide.shapes.add_textbox(Cm(1), ppt.slide_height - Cm(1.5), ppt.slide_width - Cm(2), Cm(1))
            footer_frame = footer.text_frame
            footer_frame.text = f"Page {i + 1} of {len(ppt.slides)}"
            footer_frame.paragraphs[0].font.size = Pt(styles['footer']['font_size'])
            footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        logging.debug("PowerPoint dictation created successfully.")
        return ppt
    except Exception as e:
        logging.error(f"Error creating PowerPoint: {e}")
        raise  # Re-raise the exception after logging

def add_content_slide(ppt, title, content):
    """Add a content slide to the PowerPoint presentation with a title and content."""
    try:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Content Slide
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        logging.debug("Content slide added successfully.")
        return slide
    except Exception as e:
        logging.error(f"Error adding content slide: {e}")
        raise  # Re-raise the exception after logging

# Function to calculate text size using Pillow
def calculate_text_size(text, font_path, font_size):
    """Calculate the width and height of the text using the specified font."""
    try:
        font = ImageFont.truetype(font_path, font_size)
        bbox = font.getbbox(text)
        width = bbox[2] - bbox[0]
        height = bbox[3] - bbox[1]
        # Add some padding
        width += 10
        height += 10
        return width, height
    except Exception as e:
        logging.error(f"Error calculating text size: {e}")
        raise  # Re-raise the exception after logging

# Function to add a text box with a border and adjusted margins
def add_text_box(slide, text, left, top, max_height, word_id, content_type):
    """Add a text box with the specified text to the slide."""
    try:
        # Calculate text size
        text_width, text_height = calculate_text_size(text, styles[content_type]['font_path'], styles[content_type]['font_size'])
        text_width = Inches(text_width / 72)  # Convert points to inches

        # Add the text box
        txBox = slide.shapes.add_textbox(left, top, text_width, max_height)
        tf = txBox.text_frame
        tf.text = text
        for paragraph in tf.paragraphs:
            paragraph.font.name = styles[content_type]['font_name']
            paragraph.font.size = Pt(styles[content_type]['font_size'])
            paragraph.font.color.rgb = RGBColor(*styles[content_type]['font_color'])

        txBox.line.color.rgb = RGBColor(*styles[content_type]['border_color'])
        txBox.line.width = Pt(styles[content_type]['border_width'])

        # Adjust text box margins
        tf.margin_left = Pt(styles[content_type]['margin_left'])
        tf.margin_right = Pt(styles[content_type]['margin_right'])
        tf.margin_top = Pt(styles[content_type]['margin_top'])
        tf.margin_bottom = Pt(styles[content_type]['margin_bottom'])

        # Disable AutoFit
        tf.word_wrap = False

        # Add the word ID above the text box if required
        if styles[content_type]['display_id']:
            id_box = slide.shapes.add_textbox(left, top - Inches(0.3), text_width, Inches(0.3))
            id_frame = id_box.text_frame
            id_frame.text = str(word_id)
            id_frame.paragraphs[0].font.size = Pt(styles['id']['font_size'])
            id_frame.paragraphs[0].font.color.rgb = RGBColor(*styles['id']['font_color'])
            id_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    except Exception as e:
        logging.error(f"Error adding text box: {e}")
        # Re-raise the exception after logging
        raise
